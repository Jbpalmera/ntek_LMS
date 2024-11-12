import pandas as pd
import docx2txt
import os
import sys
from pptx import Presentation
from dotenv import load_dotenv
from PyPDF2 import PdfReader
from langchain.text_splitter import CharacterTextSplitter
from langchain_openai import OpenAIEmbeddings, ChatOpenAI
from langchain_community.vectorstores import FAISS
from langchain.memory import ConversationBufferMemory
from langchain.chains import ConversationalRetrievalChain
from langchain_core.messages import HumanMessage, SystemMessage
from langchain.chains.conversation.memory import ConversationBufferWindowMemory

load_dotenv()

def get_text_from_file(file_path):
    file_extension = os.path.splitext(file_path)[1].lower()
    
    if file_extension == '.pdf':
        with open(file_path, 'rb') as f:
            pdf_reader = PdfReader(f)
            text = ""
            for page in pdf_reader.pages:
                text += page.extract_text()
        return text
    
    elif file_extension == '.txt':
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()
        
    elif file_extension == '.csv':
        df = pd.read_csv(file_path)
        return '\n'.join(df.astype(str).values.flatten())
    
    elif file_extension == '.docx':
        text = docx2txt.process(file_path)
        return text
    
    elif file_extension == '.pptx':
        presentation = Presentation(file_path)
        text = ""
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    text += shape.text + '\n'
        return text
    
    else:
        return False

def get_text_chunks(text):
    text_splitter = CharacterTextSplitter(
        separator="\n",
        chunk_size=1000,
        chunk_overlap=200,
        length_function=len
    )
    chunks = text_splitter.split_text(text)
    return chunks

def get_vectorstore(text_chunks):
    embeddings = OpenAIEmbeddings()
    vectorstore = FAISS.from_texts(texts=text_chunks, embedding=embeddings)
    return vectorstore

def get_conversation_chain(vectorstore):
    llm = ChatOpenAI()
    memory = ConversationBufferMemory(
        memory_key='chat_history', return_messages=True)
    return ConversationalRetrievalChain.from_llm(
        llm=llm,
        retriever=vectorstore.as_retriever(),
        memory=memory
    )

def get_chatopenai_response(text, openai_api_key):
    llm = ChatOpenAI(
        openai_api_key=openai_api_key,
        temperature=1,
        model_name="gpt-3.5-turbo"
    )

    messages = [
        SystemMessage(
            content="I am your Customer Care Support."
        ),
        HumanMessage(
            content=f"{text}"
        ),
    ]

    response = llm.invoke(messages)
    return response

def get_file_contexts():
    folder_path = './files'
    all_texts = []

    for filename in os.listdir(folder_path):
        file_path = os.path.join(folder_path, filename)
        
        file_text = get_text_from_file(file_path)
        if file_text is not False:
            all_texts.append(file_text)

    return all_texts 
           
def getData():
    if len(sys.argv) > 1:
        msg = sys.argv[1]
        return msg
    else:
        return "No question provided"
msg = getData()

def main():
    load_dotenv()  # Load environment variables from .env file
    openai_api_key = os.getenv("OPENAI_API_KEY")

    file_contexts = get_file_contexts()
    raw_text = '\n'.join(file_contexts)

    text_chunks = get_text_chunks(raw_text)
    vectorstore = get_vectorstore(text_chunks)
    

    while True:
        conversation_chain = get_conversation_chain(vectorstore)
        user_question = msg 
        response = conversation_chain({'question': user_question})
        chat_history = response['chat_history']
        if chat_history:
            ai_message = chat_history[1].content
            print(ai_message)
            print()
        else:
            # If no relevant response found in document context, use OpenAI
            openai_response = get_chatopenai_response(user_question, openai_api_key)
            print(openai_response.content)
            print()
        break

if __name__ == '__main__':
    main()
